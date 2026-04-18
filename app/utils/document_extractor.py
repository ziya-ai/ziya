"""
Document extraction utilities for PDF, DOC/DOCX, and XLS/XLSX files.

This module provides functions to extract text content from various document formats
so they can be meaningfully included in the context for the LLM.
"""

import os
import io
import base64
from typing import Optional, Dict, Any, List
import logging
from functools import lru_cache

from app.utils.logging_utils import logger

# Track which libraries are available
_LIBRARIES_CHECKED = False
_AVAILABLE_LIBRARIES = {
    'pypdf2': False,
    'pdfplumber': False,
    'python_docx': False,
    'openpyxl': False,
    'xlrd': False,
    'pandas': False,
    'python_pptx': False
}

# Cache for extracted document content
# Key: (file_path, mtime), Value: extracted text
_DOCUMENT_CACHE: Dict[tuple, str] = {}
_CACHE_MAX_SIZE = 100  # Maximum number of documents to cache

def _check_libraries():
    """Check which document processing libraries are available."""
    global _LIBRARIES_CHECKED, _AVAILABLE_LIBRARIES
    
    if _LIBRARIES_CHECKED:
        return
    
    # Check PDF libraries
    try:
        import pypdf
        _AVAILABLE_LIBRARIES['pypdf'] = True
    except ImportError:
        pass
    
    try:
        import pdfplumber
        _AVAILABLE_LIBRARIES['pdfplumber'] = True
    except ImportError:
        pass
    
    # Check Word document libraries
    try:
        import docx
        _AVAILABLE_LIBRARIES['python_docx'] = True
    except ImportError:
        pass
    
    # Check Excel libraries
    try:
        import openpyxl
        _AVAILABLE_LIBRARIES['openpyxl'] = True
    except ImportError:
        pass
    
    try:
        import xlrd
        _AVAILABLE_LIBRARIES['xlrd'] = True
    except ImportError:
        pass
    
    try:
        import pandas
        _AVAILABLE_LIBRARIES['pandas'] = True
    except ImportError:
        pass
    
    try:
        import pptx
        _AVAILABLE_LIBRARIES['python_pptx'] = True
    except ImportError:
        pass
    
    _LIBRARIES_CHECKED = True
    
    # Log available libraries
    available = [lib for lib, avail in _AVAILABLE_LIBRARIES.items() if avail]
    if available:
        logger.info(f"Document extraction libraries available: {', '.join(available)}")
    else:
        logger.warning("No document extraction libraries found. Install with: pip install pypdf pdfplumber python-docx openpyxl pandas python-pptx")

def is_document_file(file_path: str) -> bool:
    """
    Check if a file is a supported document type.
    
    Args:
        file_path: Path to the file
        
    Returns:
        True if the file is a supported document type
    """
    _check_libraries()
    
    # If no libraries are available, don't treat any files as documents
    if not any(_AVAILABLE_LIBRARIES.values()):
        return False
    
    if not os.path.exists(file_path):
        return False
    
    ext = os.path.splitext(file_path)[1].lower()
    supported_extensions = {'.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}
    
    return ext in supported_extensions

def is_tool_backed_file(file_path: str) -> bool:
    """
    Check if a file has specialized tool support (like pcap files).
    These files should be marked with -1 tokens to indicate tool availability.
    """
    ext = os.path.splitext(file_path)[1].lower()
    tool_backed_extensions = {'.pcap', '.pcapng', '.cap', '.dmp'}
    return ext in tool_backed_extensions

def extract_pdf_text(file_path: str) -> Optional[str]:
    """
    Extract text from a PDF file.
    
    Args:
        file_path: Path to the PDF file
        
    Returns:
        Extracted text or None if extraction failed
    """
    _check_libraries()
    
    # Try pdfplumber first (better text extraction)
    if _AVAILABLE_LIBRARIES['pdfplumber']:
        try:
            logger.debug(f"Using pdfplumber to extract from: {file_path}")
            import pdfplumber
            
            text_content = []
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            
            if text_content:
                logger.debug(f"pdfplumber: extracted text from {len(text_content)}/{page_count} pages")
            else:
                logger.warning(f"pdfplumber: opened {file_path} ({page_count} pages) but no text found — may be a scanned image")
            return '\n\n'.join(text_content) if text_content else None
            
        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path}: {e}")
    
    # Fallback to pypdf
    if _AVAILABLE_LIBRARIES['pypdf']:
        try:
            import pypdf
            
            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            
            return '\n\n'.join(text_content) if text_content else None
            
        except Exception as e:
            logger.warning(f"pypdf failed for {file_path}: {e}")
    
    logger.error(f"No PDF libraries available to extract text from {file_path}")
    return None


def extract_pdf_page_images(file_path: str, max_pages: int = 20, max_edge: int = 1568) -> Optional[List[Dict[str, Any]]]:
    """
    Render PDF pages as images.  Useful for scanned PDFs that contain no
    extractable text — the page images can be sent to a vision-capable model.

    Uses pypdfium2 (a pdfplumber dependency) to render each page.

    Args:
        file_path: Path to the PDF file.
        max_pages: Maximum number of pages to render.
        max_edge:  Maximum pixel dimension on the longer edge.

    Returns:
        A list of dicts ``{data, mediaType, page, width, height}`` where
        *data* is a base64-encoded JPEG string, or ``None`` on failure.
    """
    try:
        import pypdfium2 as pdfium
        from PIL import Image as PILImage
    except ImportError:
        logger.warning("pypdfium2 or Pillow not available — cannot render PDF pages as images")
        return None

    try:
        pdf = pdfium.PdfDocument(file_path)
    except Exception as e:
        logger.warning(f"pypdfium2 could not open {file_path}: {e}")
        return None

    page_count = len(pdf)
    if page_count == 0:
        pdf.close()
        return None

    pages_to_render = min(page_count, max_pages)
    logger.info(f"Rendering {pages_to_render}/{page_count} PDF pages as images from {file_path}")

    images: List[Dict[str, Any]] = []
    for i in range(pages_to_render):
        try:
            page = pdf[i]
            # Render at 150 DPI (good balance of quality vs size)
            bitmap = page.render(scale=150 / 72)
            pil_image = bitmap.to_pil()

            # Resize if larger than max_edge
            w, h = pil_image.size
            if max(w, h) > max_edge:
                scale = max_edge / max(w, h)
                pil_image = pil_image.resize(
                    (round(w * scale), round(h * scale)),
                    PILImage.LANCZOS,
                )
                w, h = pil_image.size

            buf = io.BytesIO()
            pil_image.save(buf, format="JPEG", quality=80)
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")

            images.append({
                "data": b64,
                "mediaType": "image/jpeg",
                "page": i + 1,
                "width": w,
                "height": h,
            })
        except Exception as e:
            logger.warning(f"Failed to render page {i + 1} of {file_path}: {e}")

    pdf.close()
    return images if images else None


def extract_docx_text(file_path: str) -> Optional[str]:
    """
    Extract text from a DOCX file.
    
    Args:
        file_path: Path to the DOCX file
        
    Returns:
        Extracted text or None if extraction failed
    """
    _check_libraries()
    
    if not _AVAILABLE_LIBRARIES['python_docx']:
        logger.error(f"python-docx library not available to extract text from {file_path}")
        return None
    
    try:
        import docx
        
        doc = docx.Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(' | '.join(row_text))
        
        return '\n\n'.join(text_content) if text_content else None
        
    except Exception as e:
        logger.error(f"Failed to extract text from DOCX {file_path}: {e}")
        return None

def extract_excel_text(file_path: str) -> Optional[str]:
    """
    Extract text from an Excel file (XLS or XLSX).
    
    Args:
        file_path: Path to the Excel file
        
    Returns:
        Extracted text or None if extraction failed
    """
    _check_libraries()
    
    # Try pandas first (handles both XLS and XLSX)
    if _AVAILABLE_LIBRARIES['pandas']:
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_content = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Convert DataFrame to text representation
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_text += df.to_string(index=False, na_rep='')
                text_content.append(sheet_text)
            
            return '\n\n'.join(text_content) if text_content else None
            
        except Exception as e:
            logger.error(f"Failed to extract text from Excel {file_path}: {e}")
            return None
    
    logger.error(f"No Excel libraries available to extract text from {file_path}")
    return None

def extract_pptx_text(file_path: str) -> Optional[str]:
    """
    Extract text from a PowerPoint file (PPT or PPTX).
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        Extracted text or None if extraction failed
    """
    _check_libraries()
    
    if not _AVAILABLE_LIBRARIES['python_pptx']:
        logger.error(f"python-pptx library not available to extract text from {file_path}")
        return None
    
    try:
        import pptx
        
        presentation = pptx.Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = f"Slide {slide_num}:"
            slide_content = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if slide_content:
                slide_text += "\n" + "\n".join(slide_content)
                text_content.append(slide_text)
        
        return '\n\n'.join(text_content) if text_content else None
        
    except Exception as e:
        logger.error(f"Failed to extract text from PowerPoint {file_path}: {e}")
        return None

def extract_document_text(file_path: str) -> Optional[str]:
    """
    Extract text from a document file based on its extension.
    Uses caching to avoid re-extracting unchanged documents.
    
    Args:
        file_path: Path to the document file
        
    Returns:
        Extracted text or None if extraction failed
    """
    global _DOCUMENT_CACHE
    
    # Check cache first
    if not os.path.exists(file_path):
        logger.error(f"Document file not found: {file_path}")
        return None
    
    try:
        mtime = os.path.getmtime(file_path)
        cache_key = (file_path, mtime)
        
        # Check if we have cached content
        if cache_key in _DOCUMENT_CACHE:
            logger.debug(f"Returning cached content for: {file_path}")
            return _DOCUMENT_CACHE[cache_key]
    except OSError as e:
        logger.warning(f"Could not get mtime for {file_path}: {e}")
        # Continue without caching
    
    # Extract the document
    extracted_text = _extract_document_text_impl(file_path)
    
    # Cache the result if extraction was successful
    if extracted_text is not None and 'cache_key' in locals():
        # Implement simple LRU by clearing oldest entries if cache is full
        if len(_DOCUMENT_CACHE) >= _CACHE_MAX_SIZE:
            # Remove oldest entry (first item in dict)
            oldest_key = next(iter(_DOCUMENT_CACHE))
            del _DOCUMENT_CACHE[oldest_key]
            logger.debug(f"Cache full, removed oldest entry: {oldest_key[0]}")
        
        _DOCUMENT_CACHE[cache_key] = extracted_text
        logger.debug(f"Cached extracted content for: {file_path} (cache size: {len(_DOCUMENT_CACHE)})")
    
    return extracted_text

def _extract_document_text_impl(file_path: str) -> Optional[str]:
    """
    Internal implementation of document text extraction.
    This is called by extract_document_text after cache check.
    
    Args:
        file_path: Path to the document file
        
    Returns:
        Extracted text or None if extraction failed
    """
    # Check if we have any document processing libraries
    _check_libraries()

    if not any(_AVAILABLE_LIBRARIES.values()):
        logger.warning(f"Cannot extract text from {file_path}: document extraction libraries not installed")
        return None
    
    logger.debug(f"Attempting to extract text from document: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        logger.debug(f"Processing PDF file: {file_path}")
        return extract_pdf_text(file_path)
    elif ext in ['.doc', '.docx']:
        logger.debug(f"Processing Word document: {file_path}")
        return extract_docx_text(file_path)
    elif ext in ['.xls', '.xlsx']:
        logger.debug(f"Processing Excel file: {file_path}")
        return extract_excel_text(file_path)
    elif ext in ['.ppt', '.pptx']:
        return extract_pptx_text(file_path)
    else:
        logger.warning(f"Unsupported document type: {ext}")
        return None
